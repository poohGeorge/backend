from . import textTranslate
import pptx
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

def translate_pptx(input_file, source_lang, target_lang):
    # Load the PowerPoint presentation
    prs = Presentation(input_file)
    output_file = "translated_pptx.pptx"
    
    # Iterate through each slide
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                # Clear original text
                original_text = shape.text_frame.text
                translated_text = textTranslate.translate_text(original_text, source_lang, target_lang)

                # Clear original text frame
                shape.text_frame.clear()

                # Create new paragraph with translated text
                new_paragraph = shape.text_frame.add_paragraph()
                new_paragraph.text = translated_text

                # Preserve original style from the first paragraph (if any)
                if shape.text_frame.paragraphs:
                    original_paragraph = shape.text_frame.paragraphs[0]
                    new_paragraph.space_after = original_paragraph.space_after
                    new_paragraph.space_before = original_paragraph.space_before
                    new_paragraph.alignment = original_paragraph.alignment
                    new_paragraph.level = original_paragraph.level

                    # Copy runs to preserve formatting
                    for run in original_paragraph.runs:
                        new_run = new_paragraph.add_run()
                        if hasattr(run.font, 'color') and run.font.color is not None:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb is not None:
                                new_run.font.color.rgb = run.font.color.rgb
                        new_run.text = translated_text
                        new_run.font.size = run.font.size
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic

    # Save the translated presentation
    prs.save(output_file)